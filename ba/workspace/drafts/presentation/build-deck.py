# -*- coding: utf-8 -*-
"""Build TOSS condensed customer deck, inheriting VNA_TOSS.pptx theme/fonts."""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

SRC = r"ba\input\VNA_TOSS.pptx"
OUT = r"ba\output\human\presentation\TOSS-trinh-bay-khach-hang-rut-gon-v1.0-2026-05-29.pptx"
ASSETS = r"ba\output\human\presentation\assets"

# Brand palette (extracted from VNA_TOSS.pptx)
DARK   = RGBColor(0x00, 0x3E, 0x54)   # brand deep teal (divider/cover bg)
ACCENT = RGBColor(0x00, 0x90, 0xB3)   # cyan accent
GOLD   = RGBColor(0xFA, 0xD1, 0x4F)   # gold highlight
TEXT   = RGBColor(0x34, 0x40, 0x54)   # primary text
TEXT2  = RGBColor(0x66, 0x70, 0x85)   # secondary text
CARD   = RGBColor(0xEA, 0xF4, 0xF8)   # light card bg
CARD2  = RGBColor(0xF2, 0xF4, 0xF7)   # grey card bg
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
BORDER = RGBColor(0xE4, 0xE7, 0xEC)
DARK2  = RGBColor(0x01, 0x50, 0x67)   # teal variant

FONT = "Inter"

prs = Presentation(SRC)
EMU = 914400
PW = prs.slide_width / EMU      # 20.0
PH = prs.slide_height / EMU     # 11.25

# remove existing slides (drop relationship so old parts are not re-serialized)
xml_slides = prs.slides._sldIdLst
for sid in list(xml_slides):
    rId = sid.get(qn('r:id'))
    try:
        prs.part.drop_rel(rId)
    except Exception:
        pass
    xml_slides.remove(sid)

blank = next(l for l in prs.slide_layouts if l.name == "Blank")

def slide():
    return prs.slides.add_slide(blank)

def notes(s, txt):
    s.notes_slide.notes_text_frame.text = txt

def _noshadow(sp):
    sp.shadow.inherit = False

def rect(s, x, y, w, h, fill=None, line=None, lw=1.0, shape=MSO_SHAPE.RECTANGLE):
    sp = s.shapes.add_shape(shape, Inches(x), Inches(y), Inches(w), Inches(h))
    if fill is None:
        sp.fill.background()
    else:
        sp.fill.solid(); sp.fill.fore_color.rgb = fill
    if line is None:
        sp.line.fill.background()
    else:
        sp.line.color.rgb = line; sp.line.width = Pt(lw)
    _noshadow(sp)
    return sp

def _fill_tf(tf, paras, anchor, lmar=0.12, tmar=0.06):
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = Inches(lmar); tf.margin_right = Inches(lmar)
    tf.margin_top = Inches(tmar); tf.margin_bottom = Inches(tmar)
    for i, p in enumerate(paras):
        para = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        para.alignment = p.get("align", PP_ALIGN.LEFT)
        if p.get("sb"): para.space_before = Pt(p["sb"])
        para.space_after = Pt(p.get("sa", 3))
        if p.get("ls"): para.line_spacing = p["ls"]
        segs = p["t"] if isinstance(p["t"], list) else [(p["t"], {})]
        for j, seg in enumerate(segs):
            txt, o = seg if isinstance(seg, tuple) else (seg, {})
            r = para.add_run(); r.text = txt
            r.font.name = FONT
            r.font.size = Pt(o.get("sz", p.get("sz", 18)))
            r.font.bold = o.get("b", p.get("b", False))
            r.font.color.rgb = o.get("c", p.get("c", TEXT))

def textbox(s, x, y, w, h, paras, anchor=MSO_ANCHOR.TOP):
    tb = s.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    _fill_tf(tb.text_frame, paras, anchor, lmar=0, tmar=0)
    return tb

def shtext(sp, paras, anchor=MSO_ANCHOR.MIDDLE, lmar=0.18):
    _fill_tf(sp.text_frame, paras, anchor, lmar=lmar)
    return sp

# ---- composite elements ----
def header(s, tag, title):
    """content-slide header: small accent tag chip + title + gold rule"""
    chip = rect(s, 1.0, 0.62, 0.13, 0.5, fill=ACCENT)  # accent tick
    textbox(s, 1.25, 0.6, 16, 0.5, [{"t": tag, "sz": 16, "b": True, "c": ACCENT}], MSO_ANCHOR.MIDDLE)
    textbox(s, 0.98, 1.12, 18, 1.0, [{"t": title, "sz": 38, "b": True, "c": TEXT}])
    rect(s, 1.0, 2.18, 2.6, 0.07, fill=GOLD)

def divider(s, num, title, en):
    rect(s, 0, 0, PW, PH, fill=DARK)
    rect(s, 1.5, 3.1, 1.4, 0.16, fill=GOLD)
    textbox(s, 1.5, 3.4, 16, 2.0, [{"t": num, "sz": 130, "b": True, "c": ACCENT}])
    textbox(s, 1.55, 6.0, 17, 1.2, [{"t": title, "sz": 50, "b": True, "c": WHITE}])
    textbox(s, 1.55, 7.1, 17, 0.8, [{"t": en, "sz": 24, "c": RGBColor(0x9F,0xC6,0xD4)}])

def card(s, x, y, w, h, paras, fill=CARD, line=None, anchor=MSO_ANCHOR.TOP, shape=MSO_SHAPE.ROUNDED_RECTANGLE, lmar=0.22):
    sp = rect(s, x, y, w, h, fill=fill, line=line, shape=shape)
    try:
        sp.adjustments[0] = 0.06
    except Exception:
        pass
    shtext(sp, paras, anchor, lmar=lmar)
    return sp

def bullets(items, sz=18, c=TEXT, sa=7, gap=None):
    out = []
    for it in items:
        if isinstance(it, tuple):
            txt, opts = it
        else:
            txt, opts = it, {}
        out.append({"t": [("•  ", {"c": ACCENT, "b": True, "sz": sz}), (txt, {"sz": opts.get("sz", sz), "c": opts.get("c", c), "b": opts.get("b", False)})], "sa": opts.get("sa", sa), "ls": 1.02})
    return out

import os
from PIL import Image
def picture(s, fname, x, y, max_w, max_h, frame=True):
    """place image fit inside (max_w,max_h) box, centered, with optional border card"""
    path = os.path.join(ASSETS, fname)
    iw, ih = Image.open(path).size
    ar = iw / ih
    w = max_w; h = w / ar
    if h > max_h:
        h = max_h; w = h * ar
    px = x + (max_w - w) / 2
    py = y + (max_h - h) / 2
    if frame:
        rect(s, px-0.18, py-0.18, w+0.36, h+0.36, fill=WHITE, line=BORDER, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    s.shapes.add_picture(path, Inches(px), Inches(py), Inches(w), Inches(h))
    return px, py, w, h

def footer(s, dark=False):
    c = RGBColor(0x9F,0xC6,0xD4) if dark else TEXT2
    textbox(s, 1.0, PH-0.55, 12, 0.4, [{"t": "TOSS — Hệ thống Tối ưu Điều hành Khai thác", "sz": 11, "c": c}])
    textbox(s, PW-4.0, PH-0.55, 3.0, 0.4, [{"t": "VTIT × VNA", "sz": 11, "c": c, "align": PP_ALIGN.RIGHT}])

# ============================ SLIDES ============================

# --- Slide 1: Cover ---
s = slide()
rect(s, 0, 0, PW, PH, fill=DARK)
rect(s, 0, 0, 0.35, PH, fill=ACCENT)
textbox(s, 1.6, 3.2, 17, 1.6, [{"t": "TOSS", "sz": 96, "b": True, "c": WHITE}])
rect(s, 1.7, 4.95, 3.2, 0.1, fill=GOLD)
textbox(s, 1.6, 5.2, 17, 1.0, [{"t": "Hệ thống Tối ưu Điều hành Khai thác", "sz": 38, "b": True, "c": RGBColor(0xEA,0xF4,0xF8)}])
textbox(s, 1.65, 6.2, 17, 0.7, [{"t": "Transformative Operations Support System", "sz": 22, "c": RGBColor(0x9F,0xC6,0xD4)}])
textbox(s, 1.6, 8.7, 17, 1.2, [
    {"t": "Trình bày giải pháp · Kick-off {{NGÀY}}", "sz": 20, "c": WHITE, "sa": 6},
    {"t": "VTIT — Công ty TNHH MTV Đầu tư Công nghệ Viettel", "sz": 18, "c": RGBColor(0x9F,0xC6,0xD4)},
])
notes(s, "Chào anh chị. Hôm nay chúng tôi trình bày giải pháp cho hệ thống TOSS — hệ thống VNA thuê trọn gói giai đoạn 2026–2031, gồm cả xây dựng và vận hành.")

# --- Slide 2: Agenda ---
s = slide()
header(s, "NỘI DUNG TRÌNH BÀY", "Buổi làm việc đi qua 5 phần")
ag = [("01", "Tổng quan"), ("02", "Các module nghiệp vụ"), ("03", "Kiến trúc & công nghệ"),
      ("04", "Triển khai & nhân sự"), ("05", "Trao đổi")]
y = 2.9
for i, (n, t) in enumerate(ag):
    cy = y + i*1.32
    c = rect(s, 1.0, cy, 1.2, 1.05, fill=DARK, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    try: c.adjustments[0] = 0.18
    except Exception: pass
    shtext(c, [{"t": n, "sz": 30, "b": True, "c": GOLD, "align": PP_ALIGN.CENTER}], MSO_ANCHOR.MIDDLE, lmar=0)
    textbox(s, 2.5, cy, 15, 1.05, [{"t": t, "sz": 26, "b": True, "c": TEXT}], MSO_ANCHOR.MIDDLE)
footer(s)
notes(s, "Buổi làm việc đi qua năm phần. Trọng tâm là nghiệp vụ và kỹ thuật. Anh chị cứ ngắt giữa chừng nếu có câu hỏi.")

# --- Slide 3: Divider 01 ---
s = slide(); divider(s, "01", "Tổng quan", "Overview")

# --- Slide 4 (deck #3): Vì sao cần TOSS ---
s = slide()
header(s, "01 · TỔNG QUAN", "Vì sao cần TOSS")
textbox(s, 1.0, 2.6, 9.5, 5.0, bullets([
    "Điều hành bay đang chuyển từ xử lý sự vụ sang quyết định dựa trên dữ liệu",
    "Ba sức ép thường trực: nhiên liệu · đúng giờ · an toàn",
    "Dữ liệu rải khắp nhiều hệ thống → quyết định chậm",
], sz=21, sa=14))
card(s, 1.0, 6.4, 9.5, 1.7, [
    {"t": "TOSS gom dữ liệu về một chỗ", "sz": 24, "b": True, "c": DARK, "sa": 4},
    {"t": "Nhìn đủ bức tranh — quyết định nhanh hơn.", "sz": 19, "c": TEXT2},
], fill=CARD, anchor=MSO_ANCHOR.MIDDLE)
# right visual: scattered systems -> TOSS
bx = 11.2
labels = ["Ops++", "AVES", "AMOS", "MO+", "LIDO"]
import math
for i, lb in enumerate(labels):
    cx = bx + (i % 3)*2.45
    cyy = 2.9 + (i // 3)*1.15
    card(s, cx, cyy, 2.2, 0.9, [{"t": lb, "sz": 16, "b": True, "c": TEXT, "align": PP_ALIGN.CENTER}], fill=CARD2, line=BORDER, anchor=MSO_ANCHOR.MIDDLE, lmar=0)
arr = rect(s, 13.4, 5.4, 2.0, 0.9, fill=ACCENT, shape=MSO_SHAPE.DOWN_ARROW)
toss = card(s, 11.4, 6.7, 6.0, 1.4, [{"t": "TOSS", "sz": 30, "b": True, "c": WHITE, "align": PP_ALIGN.CENTER}], fill=DARK, anchor=MSO_ANCHOR.MIDDLE, lmar=0)
footer(s)
notes(s, "Mỗi ngày trực ban phải cân ba thứ cùng lúc — chi phí nhiên liệu, giữ đúng giờ và an toàn — trong khi dữ liệu lại nằm rải ở hàng chục hệ thống. TOSS sinh ra để gom các luồng đó về một nơi, giúp người điều hành nhìn đủ bức tranh và quyết nhanh hơn.")

# --- Slide 5 (#4): TOSS là gì, 3 việc ---
s = slide()
header(s, "01 · TỔNG QUAN", "TOSS là gì — làm ba việc chính")
textbox(s, 1.0, 2.55, 18, 0.8, [{"t": [("Một lớp hội tụ trên các hệ thống VNA đang dùng ", {"sz": 20, "c": TEXT}),
        ("(Ops++, AVES, AMOS, MO+, LIDO…)", {"sz": 20, "c": ACCENT, "b": True})]}])
three = [
    ("1", "Điều hành chuyến bay", "Gom mọi thông tin chuyến bay, cảnh báo bất thường sớm."),
    ("2", "Tài liệu chuyến bay", "Số hóa cả bộ hồ sơ, đồng bộ trạng thái với MO+."),
    ("3", "Báo cáo & tối ưu", "Một nguồn số liệu, có khuyến nghị mang dầu (Tankering)."),
]
cw = 5.7; gap = 0.55; x0 = 1.0; yy = 3.7; ch = 4.2
for i, (n, t, d) in enumerate(three):
    cx = x0 + i*(cw+gap)
    c = card(s, cx, yy, cw, ch, [], fill=WHITE, line=BORDER, anchor=MSO_ANCHOR.TOP)
    badge = rect(s, cx+0.4, yy+0.45, 1.0, 1.0, fill=DARK, shape=MSO_SHAPE.OVAL)
    shtext(badge, [{"t": n, "sz": 30, "b": True, "c": GOLD, "align": PP_ALIGN.CENTER}], MSO_ANCHOR.MIDDLE, lmar=0)
    textbox(s, cx+0.45, yy+1.7, cw-0.9, 1.0, [{"t": t, "sz": 23, "b": True, "c": DARK}])
    textbox(s, cx+0.45, yy+2.55, cw-0.9, 1.4, [{"t": d, "sz": 18, "c": TEXT2, "ls": 1.05}])
footer(s)
notes(s, "TOSS không thay thế hệ thống nào của VNA, mà đứng trên chúng như một lớp hội tụ. Nó làm ba việc: điều hành chuyến bay với cảnh báo sớm, quản lý tài liệu đồng bộ thẳng với MO+, và báo cáo tối ưu từ một nguồn số liệu duy nhất — trong đó có cả khuyến nghị mang dầu giúp tiết kiệm chi phí.")

# --- Slide 6 (#5): Ai dùng & nền tảng ---
s = slide()
header(s, "01 · TỔNG QUAN", "Ai dùng & nền tảng")
textbox(s, 1.0, 2.5, 18, 0.5, [{"t": "NGƯỜI DÙNG", "sz": 15, "b": True, "c": ACCENT}])
users = ["Trực ban", "Điều phái", "CLC hàng hóa & hành khách", "Đại diện sân bay", "Quản trị hệ thống"]
uw = 3.5; ux = 1.0
for i, u in enumerate(users):
    card(s, ux + i*(uw+0.25), 3.05, uw, 1.0, [{"t": u, "sz": 16, "b": True, "c": TEXT, "align": PP_ALIGN.CENTER}], fill=CARD, anchor=MSO_ANCHOR.MIDDLE, lmar=0.1)
textbox(s, 1.0, 4.5, 18, 0.5, [{"t": "DỰNG TRÊN 5 NỀN TẢNG", "sz": 15, "b": True, "c": ACCENT}])
plats = ["Quản lý người dùng tập trung", "Dữ liệu chủ (Master Data)", "Cổng API", "Nền tảng dữ liệu lớn", "Giám sát 24/7"]
for i, p in enumerate(plats):
    card(s, ux + i*(uw+0.25), 5.05, uw, 1.1, [{"t": p, "sz": 16, "b": True, "c": DARK, "align": PP_ALIGN.CENTER}], fill=CARD2, line=BORDER, anchor=MSO_ANCHOR.MIDDLE, lmar=0.1)
comp = card(s, 1.0, 6.9, 18, 1.4, [
    {"t": [("Tuân thủ NĐ 53/2022 & NĐ 13/2023   ·   Mã nguồn và dữ liệu thuộc VNA   ·   SLA ≥ 99%", {"sz": 21, "b": True, "c": WHITE})], "align": PP_ALIGN.CENTER},
], fill=DARK, anchor=MSO_ANCHOR.MIDDLE, lmar=0)
footer(s)
notes(s, "Sáu nhóm người dùng chính, mỗi nhóm một góc nhìn và bộ quyền riêng. Hệ thống dựng trên năm nền tảng — quản lý người dùng, dữ liệu chủ, cổng API, nền tảng dữ liệu lớn và giám sát tập trung. Hai điểm anh chị quan tâm: tuân thủ Nghị định 53 và 13, và mã nguồn lẫn dữ liệu đều thuộc VNA.")

# --- Slide 7 (#6): Divider 02 ---
s = slide(); divider(s, "02", "Các module nghiệp vụ", "Business modules")

# --- Slide 8 (#6 content): Năm phân hệ ---
s = slide()
header(s, "02 · NGHIỆP VỤ", "Năm phân hệ nghiệp vụ")
phs = [("01", "Điều hành chuyến bay", "phân hệ lõi", True),
       ("02", "Tài liệu chuyến bay", "", False),
       ("03", "Báo cáo & tối ưu", "", False),
       ("04", "Quản lý danh mục", "nguồn dữ liệu chuẩn", False),
       ("05", "Quản trị hệ thống", "", False)]
cw = 3.5; x0 = 1.0; yy = 3.2; ch = 4.6
for i, (n, t, sub, hot) in enumerate(phs):
    cx = x0 + i*(cw+0.25)
    fill = DARK if hot else WHITE
    tc = WHITE if hot else TEXT
    nc = GOLD if hot else ACCENT
    c = card(s, cx, yy, cw, ch, [], fill=fill, line=(None if hot else BORDER))
    textbox(s, cx+0.3, yy+0.5, cw-0.6, 1.0, [{"t": n, "sz": 40, "b": True, "c": nc}])
    textbox(s, cx+0.3, yy+1.8, cw-0.6, 1.8, [{"t": t, "sz": 21, "b": True, "c": tc, "ls": 1.05}])
    if sub:
        textbox(s, cx+0.3, yy+3.6, cw-0.6, 0.7, [{"t": sub, "sz": 15, "c": (GOLD if hot else TEXT2)}])
footer(s)
notes(s, "Năm phân hệ. Số 01 là lõi nơi dữ liệu hội tụ, số 04 là nguồn dữ liệu chuẩn cho cả hệ thống. Tôi đi nhanh qua từng cái.")

# --- Slide 9 (#7): Điều hành chuyến bay ---
s = slide()
header(s, "02 · PHÂN HỆ 1", "Điều hành chuyến bay")
textbox(s, 1.0, 2.6, 11, 5.5, bullets([
    "Một màn hình, đủ thông tin: kế hoạch → dự kiến → thực tế",
    ("Cảnh báo bằng màu: hỏng hóc, VIP, chậm chuyến, thiếu phép bay…", {"b": False}),
    "Tàu bay · tổ bay · tải trọng · nhiên liệu · sân bay theo từng chuyến",
    "Nhúng Flight Radar24, điều hành nhiều hãng cùng lúc",
], sz=21, sa=16))
# mock list visual
mx, my = 12.6, 2.7
card(s, mx, my, 6.4, 5.6, [], fill=WHITE, line=BORDER, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
rect(s, mx+0.3, my+0.35, 5.8, 0.7, fill=DARK, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
textbox(s, mx+0.5, my+0.35, 5.4, 0.7, [{"t": "Danh sách chuyến bay", "sz": 16, "b": True, "c": WHITE}], MSO_ANCHOR.MIDDLE)
rowcolors = [CARD2, WHITE, CARD2, WHITE, CARD2]
flags = [GOLD, ACCENT, RGBColor(0xC0,0x50,0x4D), GOLD, ACCENT]
fl = ["VN201  HAN-SGN   delay", "VN255  SGN-DAD   OK", "VN031  HAN-NRT   MEL", "VN611  SGN-BKK   VIP", "VN160  DAD-HAN   OK"]
for i in range(5):
    ry = my+1.25 + i*0.82
    rect(s, mx+0.3, ry, 5.8, 0.72, fill=rowcolors[i], line=BORDER)
    rect(s, mx+0.3, ry, 0.14, 0.72, fill=flags[i])
    textbox(s, mx+0.6, ry, 5.4, 0.72, [{"t": fl[i], "sz": 14, "c": TEXT}], MSO_ANCHOR.MIDDLE)
footer(s)
notes(s, "Đây là màn hình trực ban ngồi cả ca. Mỗi chuyến hiện đủ thông tin từ giờ kế hoạch đến giờ thực tế. Điểm hay là cảnh báo bằng màu — hỏng hóc, khách VIP, chậm chuyến, thiếu phép bay đều bật lên ngay để không phải lục tìm.")

# --- Slide 10 (#8): Tài liệu & báo cáo ---
s = slide()
header(s, "02 · PHÂN HỆ 2 & 3", "Tài liệu & báo cáo")
c1 = card(s, 1.0, 2.7, 8.7, 5.4, [], fill=WHITE, line=BORDER)
rect(s, 1.0, 2.7, 8.7, 0.9, fill=ACCENT, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
textbox(s, 1.4, 2.7, 8, 0.9, [{"t": "Tài liệu chuyến bay", "sz": 22, "b": True, "c": WHITE}], MSO_ANCHOR.MIDDLE)
textbox(s, 1.4, 3.9, 8.0, 4.0, bullets([
    "Số hóa cả bộ hồ sơ, quản lý phiên bản",
    "Xác nhận điện tử, có Web Mobile cho tổ bay",
    "Đồng bộ trạng thái với MO+",
    "Luôn có kênh dự phòng khi tích hợp trục trặc",
], sz=18, sa=11))
c2 = card(s, 10.3, 2.7, 8.7, 5.4, [], fill=WHITE, line=BORDER)
rect(s, 10.3, 2.7, 8.7, 0.9, fill=DARK, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
textbox(s, 10.7, 2.7, 8, 0.9, [{"t": "Báo cáo & đánh giá tối ưu", "sz": 22, "b": True, "c": WHITE}], MSO_ANCHOR.MIDDLE)
textbox(s, 10.7, 3.9, 8.0, 4.0, bullets([
    "Một nguồn số liệu, lọc linh hoạt",
    "Xuất PDF / Excel / Docx",
    "Mạnh nhất: nhiên liệu, đúng giờ (OTP)",
    "Đánh giá hiệu quả từng chuyến bay",
], sz=18, sa=11))
footer(s)
notes(s, "Gộp hai phân hệ liên quan. Phần tài liệu biến cả tập hồ sơ giấy thành quy trình số, tổ bay xác nhận ngay trên điện thoại, và luôn có kênh dự phòng để không tắc trước giờ bay. Phần báo cáo, vì số liệu đã sạch và về một mối nên rất đáng tin — mạnh nhất là phân tích nhiên liệu và đúng giờ, đúng cái lãnh đạo cần.")

# --- Slide 11 (#9): Danh mục & quản trị ---
s = slide()
header(s, "02 · PHÂN HỆ 4 & 5", "Danh mục & quản trị")
card(s, 1.0, 2.7, 8.7, 5.4, [], fill=CARD, line=None)
textbox(s, 1.4, 3.0, 8, 0.7, [{"t": "Quản lý danh mục", "sz": 22, "b": True, "c": DARK}])
textbox(s, 1.4, 3.8, 8.0, 4.0, bullets([
    "Một nguồn sự thật duy nhất",
    "Có phiên bản & lịch sử thay đổi",
    "Đồng bộ MEL/CDL, defect từ AMOS",
    "Quy tắc Tankering theo giá nhiên liệu",
], sz=18, sa=11))
card(s, 10.3, 2.7, 8.7, 5.4, [], fill=CARD2, line=None)
textbox(s, 10.7, 3.0, 8, 0.7, [{"t": "Quản trị hệ thống", "sz": 22, "b": True, "c": DARK}])
textbox(s, 10.7, 3.8, 8.0, 4.0, bullets([
    "Phân quyền tới từng API",
    "Đổi quyền có hiệu lực ngay",
    "MFA, SSO, nối LDAP/AD",
    "Ghi vết toàn bộ thao tác (audit)",
], sz=18, sa=11))
footer(s)
notes(s, "Phân hệ danh mục giữ một nguồn sự thật duy nhất, tránh mỗi nơi một bản dữ liệu lệch nhau. Phân hệ quản trị lo bảo mật — phân quyền chi tiết tới từng API, và điểm khác biệt là đổi quyền là chặn được ngay, không phải chờ token hết hạn.")

# --- Slide 12 (#9b): Divider 03 ---
s = slide(); divider(s, "03", "Kiến trúc & công nghệ", "Architecture & technology")

# --- Slide 13 (#10): Kiến trúc tổng thể (sơ đồ thật - Hình 2) ---
s = slide()
header(s, "03 · KỸ THUẬT", "Kiến trúc tổng thể")
picture(s, "kien-truc-tong-the.png", 0.9, 2.5, 12.9, 6.7)
textbox(s, 14.2, 2.85, 5.2, 0.6, [{"t": "ĐIỂM CHÍNH", "sz": 15, "b": True, "c": ACCENT}])
textbox(s, 14.2, 3.5, 5.2, 5.2, bullets([
    "Microservices · cloud-native (CNCF)",
    "Mọi truy cập qua một Gateway chung",
    "Tách dữ liệu vận hành & phân tích",
    "Chừa sẵn lớp AI/ML cho tương lai",
], sz=17, sa=16))
footer(s)
notes(s, "Đây là sơ đồ kiến trúc tổng thể. Người dùng vào từ trên cùng, qua các kênh Web/Mobile, rồi tất cả đi qua một Gateway chung. Khối ứng dụng chính là năm phân hệ nghiệp vụ ta vừa xem. Bên phải có chừa sẵn lớp AI/ML. Hai điểm đáng chú ý: mọi thứ tích hợp qua một cổng chung nên dễ kiểm soát, và dữ liệu vận hành tách khỏi dữ liệu phân tích để report không làm chậm điều hành.")

# --- Slide 13b: Kiến trúc triển khai (sơ đồ thật - Hình 4) ---
s = slide()
header(s, "03 · KỸ THUẬT", "Kiến trúc triển khai")
picture(s, "kien-truc-trien-khai.png", 1.0, 2.45, 7.2, 8.1)
textbox(s, 9.0, 3.0, 10.0, 0.6, [{"t": "PHÂN VÙNG & THÀNH PHẦN", "sz": 15, "b": True, "c": ACCENT}])
textbox(s, 9.0, 3.7, 10.0, 6.0, bullets([
    "Ba vùng tách biệt: ngoài · DMZ · nội bộ",
    "Cổng vào: NGINX (HA) → API Gateway (APISIX) → Keycloak",
    "Tích hợp: MinIO (file) · Redis (cache) · Kafka (sự kiện)",
    "Dữ liệu: mỗi miền một DB, Primary–Replica qua ProxySQL",
    "Giám sát: Loki · Tempo · Mimir · Grafana",
    "CI/CD: GitLab · ArgoCD",
], sz=18, sa=15))
footer(s)
notes(s, "Đây là sơ đồ triển khai chi tiết, cho thấy độ chín của giải pháp. Ba vùng tách biệt rõ — ngoài, DMZ và nội bộ. Người dùng đi qua tường lửa, NGINX, rồi cổng API và dịch vụ xác thực Keycloak. Khối tích hợp gồm MinIO, Redis, Kafka. Dữ liệu chia theo miền, mỗi miền một cơ sở dữ liệu chạy một ghi nhiều đọc. Bên cạnh là giám sát và CI/CD. Toàn bộ là công nghệ mã nguồn mở Viettel đã vận hành thật.")

# --- Slide 14 (#11): Công nghệ, bảo mật, vận hành ---
s = slide()
header(s, "03 · KỸ THUẬT", "Công nghệ · Bảo mật · Vận hành")
cols = [
    ("Công nghệ mở, phổ biến", ["Kubernetes, Kafka", "Keycloak, Redis", "Angular, Flutter", "Viettel đã chạy thật"]),
    ("Bảo mật nhiều lớp", ["Mã hóa đường truyền", "Phân quyền theo vai trò", "Chống tấn công (WAF)", "Đạt ATTT cấp độ 3"]),
    ("Vận hành 24/7", ["Giám sát 3 chiều", "Cảnh báo sớm", "SLA ≥ 99%", "Khôi phục đúng mốc thời gian"]),
]
cw = 5.9; x0 = 1.0; yy = 2.8; ch = 5.3
heads = [ACCENT, DARK, DARK2]
for i, (t, items) in enumerate(cols):
    cx = x0 + i*(cw+0.3)
    card(s, cx, yy, cw, ch, [], fill=WHITE, line=BORDER)
    rect(s, cx, yy, cw, 1.0, fill=heads[i], shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    textbox(s, cx+0.35, yy, cw-0.7, 1.0, [{"t": t, "sz": 20, "b": True, "c": WHITE}], MSO_ANCHOR.MIDDLE)
    textbox(s, cx+0.4, yy+1.3, cw-0.8, 3.8, bullets(items, sz=18, sa=13))
footer(s)
notes(s, "Gộp ba mảng kỹ thuật còn lại. Toàn bộ công nghệ là mã nguồn mở, phổ biến, Viettel đã chạy thật nhiều dự án nên rủi ro thấp. Bảo mật làm nhiều lớp, đạt an toàn thông tin cấp độ 3. Vận hành trực 24/7 với giám sát chủ động, phát hiện trục trặc trước khi người dùng thấy, cam kết sẵn sàng tối thiểu 99%, và có sao lưu để khôi phục về đúng thời điểm trước sự cố.")

# --- Slide 15 (#12): Định cỡ & tuân thủ ---
s = slide()
header(s, "03 · KỸ THUẬT", "Định cỡ & tuân thủ")
textbox(s, 1.0, 2.5, 18, 0.6, [{"t": "ĐÁP ỨNG ĐÚNG YÊU CẦU HỒ SƠ VNA", "sz": 15, "b": True, "c": ACCENT}])
stats = [("200", "người dùng đồng thời"), ("500", "request / giây"), ("5 TB", "dữ liệu hệ thống")]
sw = 5.9
for i, (n, d) in enumerate(stats):
    cx = 1.0 + i*(sw+0.3)
    card(s, cx, 3.1, sw, 2.1, [], fill=CARD)
    textbox(s, cx+0.4, 3.25, sw-0.8, 1.1, [{"t": n, "sz": 46, "b": True, "c": DARK}])
    textbox(s, cx+0.45, 4.4, sw-0.8, 0.7, [{"t": d, "sz": 18, "c": TEXT2}])
# resources row
textbox(s, 1.0, 5.55, 18, 0.6, [{"t": "TÀI NGUYÊN ĐỀ XUẤT", "sz": 15, "b": True, "c": ACCENT}])
res = [("311", "vCPU"), ("964 GB", "RAM"), ("17.200 GB", "Lưu trữ"), ("27", "Máy ảo")]
rw = 4.4
for i, (n, d) in enumerate(res):
    cx = 1.0 + i*(rw+0.13)
    card(s, cx, 6.15, rw, 1.5, [
        {"t": n, "sz": 28, "b": True, "c": DARK, "align": PP_ALIGN.CENTER, "sa": 2},
        {"t": d, "sz": 15, "c": TEXT2, "align": PP_ALIGN.CENTER},
    ], fill=CARD2, line=BORDER, anchor=MSO_ANCHOR.MIDDLE, lmar=0)
card(s, 1.0, 7.95, 18, 1.0, [{"t": [("Tuân thủ:  NĐ 53/2022  ·  NĐ 13/2023  ·  An toàn thông tin cấp độ 3  ·  Có dư cho dự phòng & mở rộng", {"sz": 18, "b": True, "c": WHITE})], "align": PP_ALIGN.CENTER}], fill=DARK, anchor=MSO_ANCHOR.MIDDLE, lmar=0)
footer(s)
notes(s, "Cấu hình bám đúng yêu cầu hồ sơ VNA — 200 người dùng đồng thời, 500 request mỗi giây, 5 TB dữ liệu — và đã tính dư cho dự phòng cùng mở rộng về sau. Chi tiết bóc tách có đầy đủ trong tài liệu kỹ thuật nếu anh chị cần xem.")

# --- Slide 16 (#12b): Divider 04 ---
s = slide(); divider(s, "04", "Triển khai & nhân sự", "Implementation & organization")

# --- Slide 17 (#13): Kế hoạch triển khai (Gantt 6 tháng) ---
s = slide()
header(s, "04 · TRIỂN KHAI", "Kế hoạch triển khai · 6 tháng (T6–T11/2026)")
textbox(s, 12.6, 1.35, 6.7, 0.9, [{"t": "Triển khai gọn trong 6 tháng, go-live T11/2026 · vận hành 5 năm sau đó", "sz": 13, "c": TEXT2, "align": PP_ALIGN.RIGHT, "ls": 1.05}])

R_VTIT = RGBColor(0x00, 0x90, 0xB3)
R_VNA  = RGBColor(0x2E, 0x9E, 0x5B)
R_BOTH = RGBColor(0xE8, 0x94, 0x3A)
res_color = {"VTIT": R_VTIT, "VNA": R_VNA, "VTIT & VNA": R_BOTH}

NAMEW = 5.6
GX0 = 1.0 + NAMEW         # grid left = 6.6
GW = 19.0 - GX0           # grid width = 12.4
NMON = 6
COLW = GW / NMON
months = ["T6", "T7", "T8", "T9", "T10", "T11"]
top = 2.6
hdr = rect(s, 1.0, top, NAMEW, 0.5, fill=DARK)
shtext(hdr, [{"t": "Hạng mục công việc", "sz": 13, "b": True, "c": WHITE}], MSO_ANCHOR.MIDDLE, lmar=0.15)
for mi, m in enumerate(months):
    mb = rect(s, GX0+mi*COLW, top, COLW, 0.5, fill=DARK, line=WHITE, lw=0.5)
    shtext(mb, [{"t": m+"/26", "sz": 12, "b": True, "c": WHITE, "align": PP_ALIGN.CENTER}], MSO_ANCHOR.MIDDLE, lmar=0)

rows = [
    ("PHASE", "Khảo sát & phân tích", None, None, None),
    ("TASK", "Khảo sát nghiệp vụ & quy trình", 0, 0, "VTIT & VNA"),
    ("TASK", "Yêu cầu tích hợp & giao diện", 0, 1, "VTIT & VNA"),
    ("PHASE", "Phát triển & tích hợp", None, None, None),
    ("TASK", "Hạ tầng & nền tảng", 0, 1, "VTIT"),
    ("TASK", "Phát triển các phân hệ lõi", 1, 3, "VTIT"),
    ("TASK", "Dữ liệu danh mục (master data)", 2, 3, "VTIT & VNA"),
    ("TASK", "Tích hợp & kiểm thử giao diện", 2, 3, "VTIT & VNA"),
    ("PHASE", "Đào tạo & nghiệm thu", None, None, None),
    ("TASK", "Đào tạo người dùng chính (key users)", 3, 3, "VTIT"),
    ("TASK", "Kiểm thử nghiệm thu (UAT)", 4, 4, "VNA"),
    ("TASK", "Đào tạo người dùng cuối", 4, 4, "VTIT & VNA"),
    ("PHASE", "Go-live & ổn định", None, None, None),
    ("TASK", "Chuyển giao vận hành (Go-live)", 5, 5, "VTIT & VNA"),
    ("TASK", "Ổn định hệ thống", 5, 5, "VTIT & VNA"),
]
rh = 0.40
ystart = top + 0.5
for mi in range(NMON+1):
    rect(s, GX0+mi*COLW, ystart, 0.012, rh*len(rows), fill=BORDER)
y = ystart
for kind, name, a, b, resn in rows:
    if kind == "PHASE":
        rect(s, 1.0, y, NAMEW+GW, rh, fill=CARD)
        textbox(s, 1.15, y, NAMEW+GW-0.3, rh, [{"t": name, "sz": 13, "b": True, "c": DARK}], MSO_ANCHOR.MIDDLE)
    else:
        rect(s, 1.0, y, NAMEW, rh, fill=WHITE, line=BORDER, lw=0.5)
        textbox(s, 1.15, y, NAMEW-0.3, rh, [{"t": name, "sz": 12, "c": TEXT}], MSO_ANCHOR.MIDDLE)
        bx = GX0 + a*COLW + 0.07
        bw = (b-a+1)*COLW - 0.14
        bar = rect(s, bx, y+0.06, bw, rh-0.12, fill=res_color[resn], shape=MSO_SHAPE.ROUNDED_RECTANGLE)
        shtext(bar, [{"t": resn, "sz": 10, "b": True, "c": WHITE, "align": PP_ALIGN.CENTER}], MSO_ANCHOR.MIDDLE, lmar=0)
    y += rh
# legend
ly = y + 0.22
lx = 1.0
for resn, col in [("VTIT", R_VTIT), ("VNA", R_VNA), ("VTIT & VNA", R_BOTH)]:
    rect(s, lx, ly, 0.42, 0.28, fill=col, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    textbox(s, lx+0.52, ly-0.07, 2.6, 0.42, [{"t": resn, "sz": 13, "c": TEXT}], MSO_ANCHOR.MIDDLE)
    lx += 3.0
textbox(s, 11.0, ly-0.07, 8.0, 0.42, [{"t": "Mốc thời gian là đề xuất — chốt lại với đội dự án", "sz": 12, "c": TEXT2, "align": PP_ALIGN.RIGHT}], MSO_ANCHOR.MIDDLE)
footer(s)
notes(s, "Đây là kế hoạch triển khai gọn trong sáu tháng, từ tháng 6 đến tháng 11 năm 2026, trình bày dạng Gantt cho các anh dễ theo dõi. Bốn giai đoạn nối tiếp: khảo sát, phát triển và tích hợp, đào tạo và nghiệm thu, rồi go-live và ổn định trong tháng 11. Màu thanh cho biết bên chịu trách nhiệm — xanh là VTIT, lục là VNA, cam là cả hai cùng làm. Sau khi go-live là giai đoạn vận hành kéo dài 5 năm. Các mốc cụ thể chúng tôi sẽ chốt lại cùng đội dự án.")

# --- Slide 18 (#14): Cơ cấu nhân sự ---
s = slide()
header(s, "05 · NHÂN SỰ", "Cơ cấu nhân sự & phối hợp")
c1 = card(s, 1.0, 2.8, 8.6, 4.6, [], fill=WHITE, line=BORDER)
rect(s, 1.0, 2.8, 8.6, 1.0, fill=DARK, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
textbox(s, 1.4, 2.8, 8, 1.0, [{"t": "VNA / TCTHK — Chủ quản", "sz": 21, "b": True, "c": WHITE}], MSO_ANCHOR.MIDDLE)
textbox(s, 1.4, 4.0, 8.0, 3.2, bullets(["Yêu cầu & nghiệp vụ", "Nghiệm thu (UAT)", "Duyệt thay đổi Production", "An toàn thông tin"], sz=18, sa=12))
c2 = card(s, 10.4, 2.8, 8.6, 4.6, [], fill=WHITE, line=BORDER)
rect(s, 10.4, 2.8, 8.6, 1.0, fill=ACCENT, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
textbox(s, 10.8, 2.8, 8, 1.0, [{"t": "VTIT — Triển khai", "sz": 21, "b": True, "c": WHITE}], MSO_ANCHOR.MIDDLE)
textbox(s, 10.8, 4.0, 8.0, 3.2, bullets(["Kiến trúc & phát triển", "DevOps & bảo mật", "Vận hành 24/7", "Bàn giao & đào tạo"], sz=18, sa=12))
card(s, 1.0, 7.6, 18, 1.3, [{"t": [("Phối hợp:  ", {"sz": 18, "c": TEXT}), ("hội đồng thay đổi chung  ·  review thứ Sáu hàng tuần  ·  phản hồi trong 6 giờ", {"sz": 18, "b": True, "c": DARK})], "align": PP_ALIGN.CENTER}], fill=CARD, anchor=MSO_ANCHOR.MIDDLE, lmar=0)
footer(s)
notes(s, "Hai phía rõ ràng: VNA chủ trì nghiệp vụ và giữ quyền duyệt ở các điểm trọng yếu, Viettel chủ trì xây dựng và vận hành với đội trực 24/7. Hai bên gặp nhau ở hội đồng thay đổi chung và review định kỳ hàng tuần. Ranh giới rõ thì phối hợp đỡ vướng. (Chi tiết nhân sự cần điền theo hợp đồng.)")

# --- Slide 18b: Năng lực đội ngũ VTIT ---
s = slide()
header(s, "05 · NHÂN SỰ", "Năng lực đội ngũ VTIT")
textbox(s, 1.0, 2.45, 18, 0.7, [{"t": "Đội ngũ chuyên gia công nghệ với kinh nghiệm thực tế trong lĩnh vực hàng không", "sz": 20, "b": True, "c": ACCENT}])
caps = [
    ("Kinh nghiệm hàng không", "Đã triển khai hệ thống cho lĩnh vực hàng không, điều hành khai thác (OCC): {{liệt kê dự án/hệ thống cụ thể}}"),
    ("Làm việc với hãng hàng không", "Đồng hành cùng {{tên hãng hàng không}} trong {{phạm vi dự án}} — hiểu quy trình khai thác thực tế"),
    ("Hợp tác nhà sản xuất / đối tác OEM", "Kinh nghiệm tích hợp với {{nhà sản xuất tàu bay / hệ thống AMOS, LIDO, MO+…}}"),
    ("Năng lực công nghệ lõi", "Microservices · cloud-native · dữ liệu lớn · bảo mật — đã chứng minh qua nhiều dự án quy mô lớn của Viettel"),
]
cw = 8.7; ch = 2.55; gx = 1.0; gy = 3.35
for i, (t, d) in enumerate(caps):
    cx = gx + (i % 2)*(cw+0.6)
    cy = gy + (i//2)*(ch+0.35)
    card(s, cx, cy, cw, ch, [], fill=WHITE, line=BORDER)
    rect(s, cx, cy, 0.18, ch, fill=ACCENT)
    textbox(s, cx+0.45, cy+0.28, cw-0.8, 0.7, [{"t": t, "sz": 19, "b": True, "c": DARK}])
    textbox(s, cx+0.45, cy+1.05, cw-0.85, ch-1.2, [{"t": d, "sz": 16, "c": TEXT2, "ls": 1.12}])
card(s, 1.0, 9.0, 18, 0.95, [{"t": [("Quy mô đội ngũ: {{số lượng}} chuyên gia  ·  Chứng chỉ: {{PMP · K8s · bảo mật…}}  ·  Đội vận hành O&M 24/7 chuyên trách", {"sz": 15, "b": True, "c": WHITE})], "align": PP_ALIGN.CENTER}], fill=DARK, anchor=MSO_ANCHOR.MIDDLE, lmar=0)
footer(s)
notes(s, "Slide này nhấn năng lực đội ngũ, tập trung vào kinh nghiệm trong ngành hàng không — đây là điểm các sếp quan tâm. Đội ngũ VTIT đã làm hệ thống cho lĩnh vực hàng không và điều hành khai thác, từng hợp tác với hãng bay nên hiểu quy trình thực tế, và có kinh nghiệm tích hợp với các hệ thống chuyên ngành như AMOS, LIDO. Về công nghệ lõi thì microservices, cloud, dữ liệu lớn và bảo mật đều đã được Viettel chứng minh ở các dự án lớn. (Các phần trong ngoặc cần điền số liệu và tên dự án thực tế trước khi trình bày.)")

# --- Slide 19 (#15): Cam kết & giá trị ---
s = slide()
rect(s, 0, 0, PW, PH, fill=DARK)
rect(s, 0, 0, 0.35, PH, fill=GOLD)
textbox(s, 1.5, 1.4, 17, 1.0, [{"t": "Cam kết & giá trị", "sz": 44, "b": True, "c": WHITE}])
rect(s, 1.55, 2.5, 2.6, 0.08, fill=GOLD)
vals = [
    "Gom dữ liệu về một mối → quyết định nhanh, đồng bộ",
    "Tiết kiệm thật: nhiên liệu · đường bay · đúng giờ",
    "Kiến trúc mở, sẵn sàng cho AI",
    "An toàn, tuân thủ, dữ liệu thuộc VNA",
    "Đồng hành vận hành 5 năm, SLA ≥ 99%",
]
yy = 3.2
for i, v in enumerate(vals):
    cy = yy + i*1.35
    chip = rect(s, 1.5, cy, 0.7, 0.7, fill=ACCENT, shape=MSO_SHAPE.OVAL)
    shtext(chip, [{"t": "✓", "sz": 22, "b": True, "c": WHITE, "align": PP_ALIGN.CENTER}], MSO_ANCHOR.MIDDLE, lmar=0)
    textbox(s, 2.6, cy, 16, 0.8, [{"t": v, "sz": 24, "b": True, "c": RGBColor(0xEA,0xF4,0xF8)}], MSO_ANCHOR.MIDDLE)
notes(s, "Tóm lại: TOSS gom dữ liệu vốn rải rác về một chỗ để điều hành nhìn đủ và quyết nhanh, mang lại tiết kiệm cụ thể, kiến trúc mở dùng được lâu dài, an toàn và tuân thủ, dữ liệu thuộc VNA. Và chúng tôi cam kết đồng hành vận hành suốt năm năm.")

# --- Slide 20 (#16): Q&A ---
s = slide()
rect(s, 0, 0, PW, PH, fill=DARK)
textbox(s, 0, 3.6, PW, 2.0, [{"t": "Q & A", "sz": 96, "b": True, "c": WHITE, "align": PP_ALIGN.CENTER}])
textbox(s, 0, 5.7, PW, 0.8, [{"t": "Trao đổi & thảo luận", "sz": 28, "c": RGBColor(0x9F,0xC6,0xD4), "align": PP_ALIGN.CENTER}])
textbox(s, 0, 9.4, PW, 0.6, [{"t": "{{tên đầu mối}}   ·   {{email}}   ·   {{điện thoại}}", "sz": 18, "c": RGBColor(0x9F,0xC6,0xD4), "align": PP_ALIGN.CENTER}])
notes(s, "Cảm ơn anh chị. Mời đặt câu hỏi. Những điểm cần làm rõ chúng tôi sẽ ghi nhận và thống nhất bước tiếp theo.")

prs.save(OUT)
print("Saved:", OUT)
print("Slides:", len(prs.slides._sldIdLst))
